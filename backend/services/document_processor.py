import re
import os
import io
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader
from PIL import Image

from config import get_settings

settings = get_settings()


class DocumentProcessor:
    """Process various document formats and extract structured data"""

    # 食品業界の専門用語
    FOOD_TERMS = {
        "issues": [
            "離水", "離水防止", "老化", "老化防止", "膨らみ", "テクスチャ",
            "食感", "粘度", "ゲル化", "乳化", "分離防止", "保形性",
            "冷凍耐性", "レトルト耐性", "加熱安定性", "酸安定性"
        ],
        "applications": [
            "パン", "PAN", "菓子", "デザート", "ヨーグルト", "総菜",
            "ソース", "タレ", "ドレッシング", "ジャム", "ゼリー",
            "グミ", "クッキー", "ケーキ", "アイス", "レトルト",
            "冷凍食品", "弁当", "惣菜"
        ],
        "ingredients": [
            "ゲル化剤", "増粘多糖類", "ペクチン", "カラギナン",
            "キサンタンガム", "ローカストビーンガム", "タマリンドガム",
            "寒天", "ゼラチン", "デキストリン", "澱粉"
        ]
    }

    def __init__(self):
        pass

    def parse_filename(self, filename: str) -> Dict[str, str]:
        """
        Parse filename according to the naming convention:
        [アプリケーション]_[課題感]_[使用原料]_[顧客名]_[試作ID]
        """
        metadata = {
            "application": None,
            "issue": None,
            "ingredient": None,
            "customer": None,
            "trial_id": None
        }

        # Remove extension
        name_without_ext = Path(filename).stem

        # Split by underscore
        parts = name_without_ext.split("_")

        if len(parts) >= 1:
            metadata["application"] = parts[0]
        if len(parts) >= 2:
            metadata["issue"] = parts[1]
        if len(parts) >= 3:
            metadata["ingredient"] = parts[2]
        if len(parts) >= 4:
            metadata["customer"] = parts[3]
        if len(parts) >= 5:
            # Try to extract ID (usually starts with ID followed by numbers)
            for part in parts[4:]:
                if re.match(r"ID?\d+", part, re.IGNORECASE):
                    metadata["trial_id"] = part
                    break
            if not metadata["trial_id"]:
                metadata["trial_id"] = parts[4]

        return metadata

    def extract_text_from_file(self, file_content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text and structured data from a file"""
        ext = Path(filename).suffix.lower()

        if ext in [".xlsx", ".xls"]:
            return self._process_excel(file_content)
        elif ext in [".docx", ".doc"]:
            return self._process_word(file_content)
        elif ext in [".pptx", ".ppt"]:
            return self._process_powerpoint(file_content)
        elif ext == ".pdf":
            return self._process_pdf(file_content)
        elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
            return self._process_image(file_content)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _process_excel(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Process Excel file and extract structured data"""
        workbook = load_workbook(io.BytesIO(file_content), data_only=True)

        all_text = []
        structured_data = {
            "sheets": [],
            "formulations": [],  # 配合表
            "processes": []  # 製造手順
        }

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = []
            sheet_data = {
                "name": sheet_name,
                "content": [],
                "tables": []
            }

            # Check if this is a formulation sheet
            is_formulation = any(term in sheet_name for term in ["配合", "検討", "試作"])

            current_table = []
            for row in sheet.iter_rows():
                row_values = []
                row_text = []

                for cell in row:
                    if cell.value is not None:
                        value = str(cell.value).strip()
                        if value:
                            row_values.append(value)
                            row_text.append(value)

                if row_values:
                    current_table.append(row_values)
                    sheet_text.append(" | ".join(row_text))

            # Store table data
            if current_table:
                sheet_data["tables"].append(current_table)

            # Extract formulation data if applicable
            if is_formulation and current_table:
                formulation = self._extract_formulation(current_table)
                if formulation:
                    formulation["sheet_name"] = sheet_name
                    structured_data["formulations"].append(formulation)

            sheet_data["content"] = sheet_text
            structured_data["sheets"].append(sheet_data)
            all_text.extend([f"[シート: {sheet_name}]"] + sheet_text)

        return "\n".join(all_text), structured_data

    def _extract_formulation(self, table_data: List[List[str]]) -> Optional[Dict[str, Any]]:
        """Extract formulation (配合) data from table"""
        formulation = {
            "ingredients": [],
            "ratios": [],
            "notes": []
        }

        for row in table_data:
            if len(row) >= 2:
                # Check if this looks like an ingredient row
                first_cell = row[0].lower()

                # Skip header rows
                if any(term in first_cell for term in ["原料", "材料", "配合"]):
                    continue

                # Try to find numeric values (ratios)
                for i, cell in enumerate(row):
                    try:
                        value = float(cell.replace("%", "").replace("g", "").strip())
                        formulation["ingredients"].append(row[0])
                        formulation["ratios"].append(value)
                        break
                    except (ValueError, AttributeError):
                        continue

        return formulation if formulation["ingredients"] else None

    def _process_word(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Process Word document"""
        doc = DocxDocument(io.BytesIO(file_content))

        all_text = []
        structured_data = {
            "paragraphs": [],
            "tables": []
        }

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text)
                structured_data["paragraphs"].append({
                    "text": para.text,
                    "style": para.style.name if para.style else None
                })

        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
                all_text.append(" | ".join(row_data))

            structured_data["tables"].append(table_data)

        return "\n".join(all_text), structured_data

    def _process_powerpoint(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Process PowerPoint presentation"""
        prs = Presentation(io.BytesIO(file_content))

        all_text = []
        structured_data = {
            "slides": []
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_data = {
                "number": slide_num,
                "content": []
            }

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                    slide_data["content"].append(shape.text)

                # Extract table data
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(row_text))
                        slide_data["content"].append(" | ".join(row_text))

            structured_data["slides"].append(slide_data)
            all_text.append(f"[スライド {slide_num}]")
            all_text.extend(slide_text)

        return "\n".join(all_text), structured_data

    def _process_pdf(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Process PDF document"""
        reader = PdfReader(io.BytesIO(file_content))

        all_text = []
        structured_data = {
            "pages": []
        }

        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""
            all_text.append(f"[ページ {page_num}]")
            all_text.append(page_text)

            structured_data["pages"].append({
                "number": page_num,
                "content": page_text
            })

        return "\n".join(all_text), structured_data

    def _process_image(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """Process image file (placeholder for OCR integration)"""
        # For now, just return metadata
        # TODO: Integrate with Azure Document Intelligence for OCR
        image = Image.open(io.BytesIO(file_content))

        metadata = {
            "format": image.format,
            "size": image.size,
            "mode": image.mode
        }

        return f"[画像ファイル: {metadata['format']}, サイズ: {metadata['size']}]", metadata

    def chunk_text(
        self,
        text: str,
        chunk_size: int = None,
        chunk_overlap: int = None
    ) -> List[str]:
        """Split text into chunks for embedding"""
        if chunk_size is None:
            chunk_size = settings.chunk_size
        if chunk_overlap is None:
            chunk_overlap = settings.chunk_overlap

        chunks = []
        start = 0

        while start < len(text):
            end = start + chunk_size

            # Try to break at a natural point (newline or period)
            if end < len(text):
                # Look for a good break point
                for break_char in ["\n\n", "\n", "。", ". "]:
                    break_pos = text.rfind(break_char, start, end)
                    if break_pos > start:
                        end = break_pos + len(break_char)
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            # Ensure forward progress: move at least by (chunk_size - chunk_overlap)
            # or to the end if we're done
            next_start = end - chunk_overlap
            if next_start <= start:
                # Prevent infinite loop: force progress
                next_start = start + max(1, chunk_size - chunk_overlap)
            
            start = next_start

        return chunks

    def extract_keywords(self, text: str) -> Dict[str, List[str]]:
        """Extract food industry specific keywords from text"""
        keywords = {
            "issues": [],
            "applications": [],
            "ingredients": []
        }

        text_lower = text.lower()

        for category, terms in self.FOOD_TERMS.items():
            for term in terms:
                if term.lower() in text_lower:
                    keywords[category].append(term)

        return keywords
